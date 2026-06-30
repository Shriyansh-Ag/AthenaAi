import pptx
from typing import BinaryIO
from app.models.document import ExtractedDocument, DocumentMetadata, ExtractedPage, ContentBlock
from app.services.extractors.base_extractor import BaseExtractor
from app.utils.language_utils import detect_language

class PptxExtractor(BaseExtractor):
    def extract(self, file_stream: BinaryIO, filename: str) -> ExtractedDocument:
        prs = pptx.Presentation(file_stream)
        
        metadata = DocumentMetadata(
            title=prs.core_properties.title or filename,
            author=prs.core_properties.author or "",
            creation_date=str(prs.core_properties.created) if prs.core_properties.created else "",
            total_pages=len(prs.slides),
            file_type="pptx"
        )
        
        pages = []
        headings = []
        all_text = ""
        
        for i, slide in enumerate(prs.slides):
            content_blocks = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if not text:
                        continue
                        
                    all_text += text + " "
                    
                    # Heuristic: title shape is a heading
                    if shape == slide.shapes.title:
                        cb = ContentBlock(type="heading", text=text, level=1)
                        headings.append(cb)
                    else:
                        cb = ContentBlock(type="paragraph", text=text)
                        
                    content_blocks.append(cb)
                    
                elif shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text_frame.text.strip() if hasattr(cell, "text_frame") else "" for cell in row.cells]
                        table_data.append(row_data)
                        
                    cb = ContentBlock(type="table", metadata={"data": table_data})
                    content_blocks.append(cb)
                    
                elif shape.shape_type == 13: # Picture
                    image = shape.image
                    image_bytes = image.blob
                    
                    import io
                    from PIL import Image
                    import pytesseract
                    
                    try:
                        img = Image.open(io.BytesIO(image_bytes))
                        ocr_text = pytesseract.image_to_string(img).strip()
                        if ocr_text:
                            all_text += ocr_text + " "
                            cb = ContentBlock(type="paragraph", text=ocr_text)
                            content_blocks.append(cb)
                    except Exception as e:
                        print(f"OCR failed for PPTX image: {e}")
                        
                    cb = ContentBlock(type="image")
                    content_blocks.append(cb)
                    
            pages.append(ExtractedPage(page_number=i+1, content=content_blocks))
            
        metadata.language = detect_language(all_text)
        
        return ExtractedDocument(
            metadata=metadata,
            pages=pages,
            headings=headings
        )
